import os
import shutil
import logging
import yaml
import sys
import argparse
import shelve
from pathlib import Path
from threading import Thread, Lock
from queue import Queue, Empty
from collections import defaultdict
import json
import re
import email
from email import policy
from email.parser import BytesParser
import csv

# Optional libraries
try:
    import docx
except ImportError:
    docx = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from odf import text, teletype
    from odf.opendocument import load as load_odf
except ImportError:
    load_odf = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import toml
except ImportError:
    toml = None

try:
    import extract_msg
except ImportError:
    extract_msg = None

import xml.etree.ElementTree as ET


def read_in_chunks(file_obj, chunk_size=1024 * 1024):
    """
    Read a file in specified chunk sizes.
    """
    while True:
        data = file_obj.read(chunk_size)
        if not data:
            break
        yield data


class Statistics:
    """
    Thread-safe statistics collector.
    """

    def __init__(self):
        self.lock = Lock()
        self.keyword_files = defaultdict(set)
        self.keyword_group_counts = defaultdict(int)
        self.keyword_counts = defaultdict(lambda: defaultdict(int))
        self.filetype_counts = defaultdict(int)
        self.files_with_multiple_keywords = 0

        # Track matched file sizes to compute total GB. 
        # A file might match multiple groups, so track globally once.
        self.all_matched_files = dict()

        # Track file sizes per group to facilitate separate stats for each group.
        self.group_matched_files = defaultdict(dict)

    def update_keyword_group(self, group):
        with self.lock:
            self.keyword_group_counts[group] += 1

    def update_keyword(self, group, keyword):
        with self.lock:
            self.keyword_counts[group][keyword] += 1

    def update_filetype(self, filetype):
        with self.lock:
            self.filetype_counts[filetype] += 1

    def increment_multiple_keywords(self):
        with self.lock:
            self.files_with_multiple_keywords += 1

    def update_keyword_file_path(self, keyword, file_path):
        with self.lock:
            self.keyword_files[keyword].add(str(file_path.resolve()))

    def record_matched_file(self, file_path, group, size):
        """
        Store the matched file size in global and group-based stats.
        """
        with self.lock:
            self.all_matched_files[file_path] = size
            # Store size under the group as well
            if file_path not in self.group_matched_files[group]:
                self.group_matched_files[group][file_path] = size

    def to_dict(self):
        with self.lock:
            # Compute the global total matched files size in GB
            total_bytes = sum(self.all_matched_files.values())
            total_gb = total_bytes / (1024 ** 3)

            return {
                "Keyword_Groups": dict(self.keyword_group_counts),
                "Keywords": {k: dict(v) for k, v in self.keyword_counts.items()},
                "FileTypes": dict(self.filetype_counts),
                "Files_With_Multiple_Keywords": self.files_with_multiple_keywords,
                "FilesContainingKeyword": {
                    keyword: sorted(paths)
                    for keyword, paths in self.keyword_files.items()
                },
                "Global_Matched_Files_Size_GB": round(total_gb, 4)
            }

    def group_to_dict(self, group):
        """
        Build a per-group dictionary containing only group-based statistics.
        """
        with self.lock:
            # Summarize the matched files for this group
            matched_files = self.group_matched_files[group]
            total_bytes = sum(matched_files.values())
            total_gb = total_bytes / (1024 ** 3)

            # Gather keyword counts for this group
            group_keywords_data = {}
            if group in self.keyword_counts:
                group_keywords_data = dict(self.keyword_counts[group])

            # Build separate group stats
            group_dict = {
                "Group_Name": group,
                "Keyword_Group_Count": self.keyword_group_counts.get(group, 0),
                "Keywords_in_Group": group_keywords_data,
                "Matched_Files_Size_GB": round(total_gb, 4),
                "Files_Containing_Group_Keywords": {}
            }

            # For readability, gather which keywords correspond to the files
            # from the global "keyword_files", but filter for this group only.
            # This is optional, as we already store it globally.
            # We'll do a minimal subset: only those keywords that belong to this group
            # from self.keyword_counts[group].
            relevant_keywords = set(group_keywords_data.keys())
            for keyword, paths in self.keyword_files.items():
                if keyword in relevant_keywords:
                    group_dict["Files_Containing_Group_Keywords"][keyword] = sorted(paths)

            return group_dict


def load_config(config_path):
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    if 'Keyword_Groups' not in config:
        logging.error("Configuration file is missing 'Keyword_Groups' section.")
        sys.exit(1)
    keyword_groups = config['Keyword_Groups']
    if not isinstance(keyword_groups, dict):
        logging.error("'Keyword_Groups' should be a dictionary.")
        sys.exit(1)
    for group, keywords in keyword_groups.items():
        if not isinstance(keywords, list):
            logging.error(f"Keywords for group '{group}' should be a list.")
            sys.exit(1)
    return keyword_groups


def setup_logging(log_file, verbose):
    logging.basicConfig(
        filename=log_file,
        filemode='w',
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )
    console = logging.StreamHandler()
    if verbose:
        console.setLevel(logging.INFO)
    else:
        console.setLevel(logging.WARNING)
    formatter = logging.Formatter('%(levelname)s - %(message)s')
    console.setFormatter(formatter)
    logging.getLogger().addHandler(console)


def extract_text_from_xml(element):
    text_content = element.text or ""
    for child in element:
        text_content += '\n' + extract_text_from_xml(child)
        if child.tail:
            text_content += '\n' + child.tail
    return text_content


def extract_text(file_path):
    """
    Extract text from various file types using chunked reading for large text-based files.
    """
    ext = file_path.suffix.lower()
    # If file has no extension, treat as .txt
    if not ext:
        ext = '.txt'
    text_content = ""

    try:
        if ext in ['.txt', '.csv', '.json', '.sql', '.conf', '.cfg', '.ini', '.toml']:
            # Use chunked reading
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                for chunk in read_in_chunks(f):
                    text_content += chunk
        elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
            if docx:
                try:
                    doc = docx.Document(file_path)
                    text_content = '\n'.join([para.text for para in doc.paragraphs])
                except Exception as e:
                    logging.error(f"Error reading .docx-like file {file_path}: {e}")
            else:
                logging.warning("docx library not installed.")
        elif ext in ['.xlsx', '.xls', '.xlsm', '.xltx', '.xltm']:
            if ext == '.xlsx' and openpyxl:
                try:
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows(values_only=True):
                            row_text = ' '.join([str(cell) for cell in row if cell is not None])
                            text_content += row_text + '\n'
                    wb.close()
                except Exception as e:
                    logging.error(f"Error reading .xlsx file {file_path}: {e}")
            else:
                if xlrd:
                    try:
                        workbook = xlrd.open_workbook(file_path, on_demand=True)
                        for sheet_name in workbook.sheet_names():
                            worksheet = workbook.sheet_by_name(sheet_name)
                            for row_i in range(worksheet.nrows):
                                text_content += ' '.join(
                                    [str(cell) for cell in worksheet.row(row_i)]
                                ) + '\n'
                    except Exception as e:
                        logging.error(f"Error reading .xls-like file {file_path}: {e}")
                else:
                    logging.error("xlrd not installed.")
        elif ext in ['.pptx', '.pptm', '.potx', '.potm']:
            if pptx:
                try:
                    prs = pptx.Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + '\n'
                except Exception as e:
                    logging.error(f"Error reading .pptx-like file {file_path}: {e}")
            else:
                logging.warning("pptx library not installed.")
        elif ext in ['.odt', '.ods', '.odp']:
            if load_odf:
                try:
                    odf_doc = load_odf(file_path)
                    if ext == '.odt':
                        allparas = odf_doc.getElementsByType(text.P)
                        text_content = '\n'.join([teletype.extractText(p) for p in allparas])
                    elif ext == '.ods':
                        sheets = odf_doc.spreadsheet.getElementsByType(text.Table)
                        for sheet in sheets:
                            rows = sheet.getElementsByType(text.TableRow)
                            for row in rows:
                                cells = row.getElementsByType(text.TableCell)
                                row_text = ' '.join([teletype.extractText(cell) for cell in cells])
                                text_content += row_text + '\n'
                    elif ext == '.odp':
                        slides = odf_doc.getElementsByType(text.Slide)
                        for slide in slides:
                            allparas = slide.getElementsByType(text.P)
                            text_content += '\n'.join([teletype.extractText(p) for p in allparas]) + '\n'
                except Exception as e:
                    logging.error(f"Error reading LibreOffice file {file_path}: {e}")
            else:
                logging.warning(f"odfpy not installed for {ext} files.")
        elif ext == '.pdf' and PyPDF2:
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text_content += extracted + '\n'
            except Exception as e:
                logging.error(f"Error reading PDF file {file_path}: {e}")
        elif ext == '.xml':
            try:
                tree = ET.parse(file_path)
                root = tree.getroot()
                text_content = extract_text_from_xml(root)
            except ET.ParseError as e:
                logging.error(f"Error parsing XML file {file_path}: {e}")
            except Exception as e:
                logging.error(f"Unexpected error processing XML file {file_path}: {e}")
        elif ext == '.eml':
            try:
                with open(file_path, 'rb') as f:
                    msg = BytesParser(policy=policy.default).parse(f)
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        text_content += part.get_content()
                    elif part.get_content_type() == 'text/html':
                        html_content = part.get_content()
                        clean_text = re.sub('<[^<]+?>', '', html_content)
                        text_content += clean_text + '\n'
            except Exception as e:
                logging.error(f"Error reading EML file {file_path}: {e}")
        elif ext == '.msg':
            if extract_msg:
                try:
                    msg = extract_msg.Message(file_path)
                    msg_sender = msg.sender
                    msg_date = msg.date
                    msg_subject = msg.subject
                    msg_body = msg.body
                    msg_text = (
                        f"Sender: {msg_sender}\nDate: {msg_date}\nSubject: {msg_subject}\n\n{msg_body}"
                    )
                    text_content = msg_text
                except Exception as e:
                    logging.error(f"Error reading MSG file {file_path}: {e}")
            else:
                logging.error("extract-msg not installed.")
        else:
            logging.warning(f"Unsupported file type: {file_path}")
    except Exception as e:
        logging.error(f"Error reading {file_path}: {e}")

    return text_content.lower()


def compile_keyword_patterns(keyword_groups):
    compiled_patterns = {}
    for group, keywords in keyword_groups.items():
        compiled_patterns[group] = []
        for keyword in keywords:
            escaped_keyword = re.escape(keyword.lower())
            pattern_str = r'\b' + escaped_keyword + r'\b'
            try:
                compiled = re.compile(pattern_str)
                compiled_patterns[group].append((keyword, compiled))
            except re.error as e:
                logging.error(f"Invalid regex pattern for keyword '{keyword}' in group '{group}': {e}")
    return compiled_patterns


def search_keywords(text_content, compiled_patterns):
    found_groups = {}
    for group, patterns in compiled_patterns.items():
        for keyword, pattern in patterns:
            if pattern.search(text_content):
                if group not in found_groups:
                    found_groups[group] = set()
                found_groups[group].add(keyword)
    return found_groups


def copy_file_to_groups(file_path, groups, output_dir):
    copied_groups = []
    for group in groups:
        group_folder = output_dir / group
        group_folder.mkdir(parents=True, exist_ok=True)
        destination = group_folder / file_path.name
        if destination.exists():
            base, ext = os.path.splitext(file_path.name)
            counter = 1
            while True:
                new_name = f"{base}_{counter}{ext}"
                new_destination = group_folder / new_name
                if not new_destination.exists():
                    destination = new_destination
                    break
                counter += 1
        try:
            shutil.copy2(file_path, destination)
            logging.info(f"Copied {file_path} to {destination}")
            copied_groups.append(group)
        except Exception as e:
            logging.error(f"Error copying {file_path} to {destination}: {e}")
    return copied_groups


def process_file(file_path, compiled_patterns, output_dir, stats, metadata_db):
    """
    Check for file changes, extract text if new or modified, search for keywords,
    copy matched files, and update stats.
    """
    try:
        mtime = file_path.stat().st_mtime
        fsize = file_path.stat().st_size
    except Exception as e:
        logging.error(f"Cannot access file stats for {file_path}: {e}")
        return

    db_key = str(file_path.resolve())
    stored = metadata_db.get(db_key)

    # If unchanged, skip
    if stored and stored == (mtime, fsize):
        logging.info(f"Skipping unchanged file: {file_path}")
        return

    logging.info(f"Processing file: {file_path}")
    text_content = extract_text(file_path)
    if text_content:
        found = search_keywords(text_content, compiled_patterns)
        if found:
            # Record matched file info for global sums and for each group
            for group, keywords in found.items():
                stats.update_keyword_group(group)
                for keyword in keywords:
                    stats.update_keyword(group, keyword)
                    stats.update_keyword_file_path(keyword, file_path)

                # Record file in matched stats with size
                stats.record_matched_file(file_path, group, fsize)

            if len(found) > 1:
                stats.increment_multiple_keywords()

            copy_file_to_groups(file_path, found.keys(), output_dir)
            for group, keywords in found.items():
                for keyword in keywords:
                    logging.info(f"Found keyword '{keyword}' in {file_path} for group '{group}'")

        # Update filetype stats
        stats.update_filetype(file_path.suffix.lower())
    else:
        logging.warning(f"No text extracted from {file_path}")

    # Mark current file as processed
    metadata_db[db_key] = (mtime, fsize)


def generate_group_statistics(stats, output_dir):
    """
    Generate separate group-based statistics files.
    Each file is named 'statistics_{group}.json'.
    """
    # For each group, build a stats dict and write it out
    all_groups = list(stats.keyword_group_counts.keys())
    for group in all_groups:
        group_dict = stats.group_to_dict(group)
        group_file = output_dir / f"statistics_{group}.json"
        try:
            with open(group_file, 'w', encoding='utf-8') as f:
                json.dump(group_dict, f, indent=4)
            logging.info(f"Group statistics written to {group_file}")
        except Exception as e:
            logging.error(f"Error writing group statistics for {group} to {group_file}: {e}")


def generate_statistics(stats, output_dir):
    """
    Generate global statistics, then group-based statistics files.
    """
    stats_data = stats.to_dict()
    stats_file = output_dir / "statistics.json"
    try:
        with open(stats_file, 'w', encoding='utf-8') as f:
            json.dump(stats_data, f, indent=4)
        logging.info(f"Statistics written to {stats_file}")
    except Exception as e:
        logging.error(f"Error writing statistics to {stats_file}: {e}")

    # Now generate separate files for each group
    generate_group_statistics(stats, output_dir)


def consumer(queue, compiled_patterns, output_path, stats, metadata_db):
    """
    Consumer thread function to process files from the queue.
    """
    while True:
        try:
            file_path = queue.get_nowait()
        except Empty:
            return
        try:
            process_file(file_path, compiled_patterns, output_path, stats, metadata_db)
        except Exception as e:
            logging.error(f"Unhandled exception in consumer for {file_path}: {e}")
        finally:
            queue.task_done()


def main(data_dir, config_path, output_dir, log_file, verbose):
    setup_logging(log_file, verbose)
    keyword_groups = load_config(config_path)
    compiled_patterns = compile_keyword_patterns(keyword_groups)
    stats = Statistics()

    data_path = Path(data_dir)
    if not data_path.exists():
        logging.error(f"Data directory does not exist: {data_dir}")
        sys.exit(1)

    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    supported_exts = [
        '.txt', '.csv', '.json', '.sql', '.conf', '.cfg',
        '.docx', '.docm', '.dotx', '.dotm',
        '.xlsx', '.xls', '.xlsm', '.xltx', '.xltm',
        '.pptx', '.pptm', '.potx', '.potm',
        '.odt', '.ods', '.odp',
        '.pdf', '.xml', '.ini', '.toml',
        '.eml', '.msg'
    ]

    file_queue = Queue()
    for root, dirs, files in os.walk(data_path):
        for file in files:
            file_path = Path(root) / file
            # Accept files without extension or in the supported list
            if not file_path.suffix or file_path.suffix.lower() in supported_exts:
                file_queue.put(file_path)
            else:
                logging.warning(f"Skipping unsupported file type: {file_path}")

    if file_queue.empty():
        logging.info("No supported files found to process.")
        return

    max_workers = min(32, os.cpu_count() + 4)

    # Use shelve as a small key-value store for metadata
    with shelve.open('metadata_store.db') as metadata_db:
        workers = []
        for _ in range(max_workers):
            t = Thread(
                target=consumer,
                args=(file_queue, compiled_patterns, output_path, stats, metadata_db)
            )
            t.start()
            workers.append(t)

        file_queue.join()

        for w in workers:
            w.join()

    generate_statistics(stats, output_path)
    logging.info("Processing completed.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Keyword Search and File Organizer")
    parser.add_argument('--data', default='./data', help="Path to the data directory (default: ./data)")
    parser.add_argument('--config', default='./config.yaml', help="Path to the config.yaml file (default: ./config.yaml)")
    parser.add_argument('--output', default='./output', help="Directory to store copied files (default: ./output)")
    parser.add_argument('--log', default='./script.log', help="Log file path (default: ./script.log)")
    parser.add_argument('--verbose', action='store_true', help="Enable verbose mode")
    args = parser.parse_args()
    main(args.data, args.config, args.output, args.log, args.verbose)
