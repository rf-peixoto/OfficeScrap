#!/usr/bin/env python3
import os
import sys
import re
import json
import csv
import shutil
import yaml
import argparse
import logging
import sqlite3
from pathlib import Path
from threading import Thread, Lock
from queue import Queue

import email
from email import policy
from email.parser import BytesParser
import xml.etree.ElementTree as ET
from collections import defaultdict

# Optional libraries
try:
    import docx
except ImportError:
    docx = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from odf import text, teletype
    from odf.opendocument import load as load_odf
except ImportError:
    load_odf = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import toml
except ImportError:
    toml = None

try:
    import extract_msg
except ImportError:
    extract_msg = None


# -------------------------
# Utilities
# -------------------------
def read_in_chunks(file_obj, chunk_size=1024 * 1024):
    """Read a file in specified chunk sizes."""
    while True:
        data = file_obj.read(chunk_size)
        if not data:
            break
        yield data


def setup_logging(log_file: Path, verbose: bool) -> None:
    logging.basicConfig(
        filename=str(log_file),
        filemode="w",
        level=logging.INFO,
        format="%(asctime)s - %(levelname)s - %(message)s",
    )
    console = logging.StreamHandler()
    console.setLevel(logging.INFO if verbose else logging.WARNING)
    console.setFormatter(logging.Formatter("%(levelname)s - %(message)s"))
    logging.getLogger().addHandler(console)


def safe_relpath(file_path: Path, root_dir: Path) -> Path:
    """
    Return a path relative to root_dir, avoiding leakage of the system absolute path.
    If file_path is not under root_dir, fall back to basename.
    """
    try:
        return file_path.resolve().relative_to(root_dir.resolve())
    except Exception:
        return Path(file_path.name)


def extract_text_from_xml(element: ET.Element) -> str:
    text_content = element.text or ""
    for child in element:
        text_content += "\n" + extract_text_from_xml(child)
        if child.tail:
            text_content += "\n" + child.tail
    return text_content


# -------------------------
# Metadata store (SQLite)
# -------------------------
class MetadataStore:
    """
    Thread-safe metadata store backed by SQLite.
    Tracks (mtime, size) by absolute file path (resolved) as key.
    """

    def __init__(self, db_path: Path):
        self.db_path = db_path
        self.lock = Lock()
        self.conn = sqlite3.connect(str(db_path), check_same_thread=False)
        self.conn.execute("PRAGMA journal_mode=WAL;")
        self.conn.execute(
            """
            CREATE TABLE IF NOT EXISTS metadata (
                path TEXT PRIMARY KEY,
                mtime REAL NOT NULL,
                size INTEGER NOT NULL
            )
            """
        )
        self.conn.commit()

    def get(self, key: str):
        with self.lock:
            cur = self.conn.execute("SELECT mtime, size FROM metadata WHERE path = ?", (key,))
            row = cur.fetchone()
            return tuple(row) if row else None

    def set(self, key: str, mtime: float, size: int) -> None:
        with self.lock:
            self.conn.execute(
                """
                INSERT INTO metadata(path, mtime, size)
                VALUES (?, ?, ?)
                ON CONFLICT(path) DO UPDATE SET
                    mtime=excluded.mtime,
                    size=excluded.size
                """,
                (key, float(mtime), int(size)),
            )
            self.conn.commit()

    def close(self) -> None:
        with self.lock:
            try:
                self.conn.commit()
            finally:
                self.conn.close()


# -------------------------
# Statistics
# -------------------------
class Statistics:
    """Thread-safe statistics collector."""

    def __init__(self):
        self.lock = Lock()
        self.keyword_files = defaultdict(set)  # keyword -> set(relpaths)
        self.keyword_group_counts = defaultdict(int)
        self.keyword_counts = defaultdict(lambda: defaultdict(int))
        self.filetype_counts = defaultdict(int)
        self.files_with_multiple_keywords = 0

        # Matched file sizes (global and per-group)
        self.all_matched_files = dict()          # relpath -> size
        self.group_matched_files = defaultdict(dict)  # group -> {relpath: size}

    def update_keyword_group(self, group: str) -> None:
        with self.lock:
            self.keyword_group_counts[group] += 1

    def update_keyword(self, group: str, keyword: str) -> None:
        with self.lock:
            self.keyword_counts[group][keyword] += 1

    def update_filetype(self, filetype: str) -> None:
        with self.lock:
            self.filetype_counts[filetype] += 1

    def increment_multiple_keywords(self) -> None:
        with self.lock:
            self.files_with_multiple_keywords += 1

    def update_keyword_file_path(self, keyword: str, rel_path: Path) -> None:
        with self.lock:
            self.keyword_files[keyword].add(str(rel_path.as_posix()))

    def record_matched_file(self, rel_path: Path, group: str, size: int) -> None:
        with self.lock:
            rp = str(rel_path.as_posix())
            self.all_matched_files[rp] = int(size)
            if rp not in self.group_matched_files[group]:
                self.group_matched_files[group][rp] = int(size)

    def to_dict(self) -> dict:
        with self.lock:
            total_bytes = sum(self.all_matched_files.values())
            total_gb = total_bytes / (1024**3)

            return {
                "Keyword_Groups": dict(self.keyword_group_counts),
                "Keywords": {k: dict(v) for k, v in self.keyword_counts.items()},
                "FileTypes": dict(self.filetype_counts),
                "Files_With_Multiple_Keywords": self.files_with_multiple_keywords,
                "FilesContainingKeyword": {
                    keyword: sorted(paths) for keyword, paths in self.keyword_files.items()
                },
                "Global_Matched_Files_Size_GB": round(total_gb, 4),
            }

    def group_to_dict(self, group: str) -> dict:
        with self.lock:
            matched_files = self.group_matched_files[group]
            total_bytes = sum(matched_files.values())
            total_gb = total_bytes / (1024**3)

            group_keywords_data = dict(self.keyword_counts.get(group, {}))
            relevant_keywords = set(group_keywords_data.keys())

            group_dict = {
                "Group_Name": group,
                "Keyword_Group_Count": self.keyword_group_counts.get(group, 0),
                "Keywords_in_Group": group_keywords_data,
                "Matched_Files_Size_GB": round(total_gb, 4),
                "Files_Containing_Group_Keywords": {},
            }

            for keyword, paths in self.keyword_files.items():
                if keyword in relevant_keywords:
                    group_dict["Files_Containing_Group_Keywords"][keyword] = sorted(paths)

            return group_dict


# -------------------------
# Config / matching
# -------------------------
def load_config(config_path: Path) -> dict:
    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    if "Keyword_Groups" not in config:
        logging.error("Configuration file is missing 'Keyword_Groups' section.")
        sys.exit(1)

    keyword_groups = config["Keyword_Groups"]
    if not isinstance(keyword_groups, dict):
        logging.error("'Keyword_Groups' should be a dictionary.")
        sys.exit(1)

    for group, keywords in keyword_groups.items():
        if not isinstance(keywords, list):
            logging.error(f"Keywords for group '{group}' should be a list.")
            sys.exit(1)

    return keyword_groups


def compile_keyword_patterns(keyword_groups: dict) -> dict:
    """
    Mantém a estratégia solicitada: \b<keyword>\b (com escape do texto do keyword).
    """
    compiled_patterns = {}
    for group, keywords in keyword_groups.items():
        compiled_patterns[group] = []
        for keyword in keywords:
            escaped_keyword = re.escape(str(keyword).lower())
            pattern_str = r"\b" + escaped_keyword + r"\b"
            try:
                compiled_patterns[group].append((keyword, re.compile(pattern_str)))
            except re.error as e:
                logging.error(f"Invalid regex pattern for keyword '{keyword}' in group '{group}': {e}")
    return compiled_patterns


def search_keywords(text_content: str, compiled_patterns: dict) -> dict:
    found_groups = {}
    for group, patterns in compiled_patterns.items():
        for keyword, pattern in patterns:
            if pattern.search(text_content):
                found_groups.setdefault(group, set()).add(keyword)
    return found_groups


# -------------------------
# Text extraction (mantida)
# -------------------------
def extract_text(file_path: Path) -> str:
    """
    Extract text from various file types using chunked reading for large text-based files.
    Mantido conforme solicitado (sem objetos embutidos).
    """
    ext = file_path.suffix.lower()
    if not ext:
        ext = ".txt"

    text_content = ""
    try:
        if ext in [".txt", ".csv", ".json", ".sql", ".conf", ".cfg", ".ini", ".toml"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                for chunk in read_in_chunks(f):
                    text_content += chunk

        elif ext in [".docx", ".docm", ".dotx", ".dotm"]:
            if docx:
                try:
                    doc = docx.Document(file_path)
                    text_content = "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    logging.error(f"Error reading .docx-like file {file_path}: {e}")
            else:
                logging.warning("docx library not installed.")

        elif ext in [".xlsx", ".xls", ".xlsm", ".xltx", ".xltm"]:
            if ext == ".xlsx" and openpyxl:
                try:
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows(values_only=True):
                            row_text = " ".join([str(cell) for cell in row if cell is not None])
                            text_content += row_text + "\n"
                    wb.close()
                except Exception as e:
                    logging.error(f"Error reading .xlsx file {file_path}: {e}")
            else:
                if xlrd:
                    try:
                        workbook = xlrd.open_workbook(file_path, on_demand=True)
                        for sheet_name in workbook.sheet_names():
                            worksheet = workbook.sheet_by_name(sheet_name)
                            for row_i in range(worksheet.nrows):
                                text_content += " ".join([str(cell) for cell in worksheet.row(row_i)]) + "\n"
                    except Exception as e:
                        logging.error(f"Error reading .xls-like file {file_path}: {e}")
                else:
                    logging.error("xlrd not installed.")

        elif ext in [".pptx", ".pptm", ".potx", ".potm"]:
            if pptx:
                try:
                    prs = pptx.Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + "\n"
                except Exception as e:
                    logging.error(f"Error reading .pptx-like file {file_path}: {e}")
            else:
                logging.warning("pptx library not installed.")

        elif ext in [".odt", ".ods", ".odp"]:
            if load_odf:
                try:
                    odf_doc = load_odf(str(file_path))
                    if ext == ".odt":
                        allparas = odf_doc.getElementsByType(text.P)
                        text_content = "\n".join([teletype.extractText(p) for p in allparas])
                    elif ext == ".ods":
                        sheets = odf_doc.spreadsheet.getElementsByType(text.Table)
                        for sheet in sheets:
                            rows = sheet.getElementsByType(text.TableRow)
                            for row in rows:
                                cells = row.getElementsByType(text.TableCell)
                                row_text = " ".join([teletype.extractText(cell) for cell in cells])
                                text_content += row_text + "\n"
                    elif ext == ".odp":
                        slides = odf_doc.getElementsByType(text.Slide)
                        for slide in slides:
                            allparas = slide.getElementsByType(text.P)
                            text_content += "\n".join([teletype.extractText(p) for p in allparas]) + "\n"
                except Exception as e:
                    logging.error(f"Error reading LibreOffice file {file_path}: {e}")
            else:
                logging.warning(f"odfpy not installed for {ext} files.")

        elif ext == ".pdf" and PyPDF2:
            try:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text_content += extracted + "\n"
            except Exception as e:
                logging.error(f"Error reading PDF file {file_path}: {e}")

        elif ext == ".xml":
            try:
                tree = ET.parse(file_path)
                root = tree.getroot()
                text_content = extract_text_from_xml(root)
            except ET.ParseError as e:
                logging.error(f"Error parsing XML file {file_path}: {e}")
            except Exception as e:
                logging.error(f"Unexpected error processing XML file {file_path}: {e}")

        elif ext == ".eml":
            try:
                with open(file_path, "rb") as f:
                    msg = BytesParser(policy=policy.default).parse(f)
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_content += part.get_content()
                    elif part.get_content_type() == "text/html":
                        html_content = part.get_content()
                        clean_text = re.sub("<[^<]+?>", "", html_content)
                        text_content += clean_text + "\n"
            except Exception as e:
                logging.error(f"Error reading EML file {file_path}: {e}")

        elif ext == ".msg":
            if extract_msg:
                try:
                    msg = extract_msg.Message(str(file_path))
                    msg_text = (
                        f"Sender: {msg.sender}\n"
                        f"Date: {msg.date}\n"
                        f"Subject: {msg.subject}\n\n"
                        f"{msg.body}"
                    )
                    text_content = msg_text
                except Exception as e:
                    logging.error(f"Error reading MSG file {file_path}: {e}")
            else:
                logging.error("extract-msg not installed.")
        else:
            logging.warning(f"Unsupported file type: {file_path}")

    except Exception as e:
        logging.error(f"Error reading {file_path}: {e}")

    return text_content.lower()


# -------------------------
# Copy logic (corrigida)
# -------------------------
def unique_destination_path(dest_path: Path) -> Path:
    """
    If dest_path exists, append _N before suffix (keeping directory structure).
    """
    if not dest_path.exists():
        return dest_path

    parent = dest_path.parent
    stem = dest_path.stem
    suffix = dest_path.suffix
    counter = 1
    while True:
        candidate = parent / f"{stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


def copy_file_to_groups(file_path: Path, rel_path: Path, groups, output_dir: Path) -> None:
    """
    Copy preserving the relative path under output_dir/<group>/<rel_path>.
    This prevents leaking system absolute path while keeping "absolute within root" structure.
    """
    for group in groups:
        group_root = output_dir / group
        destination = group_root / rel_path
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination = unique_destination_path(destination)

        try:
            shutil.copy2(file_path, destination)
            logging.info(f"Copied {file_path} to {destination}")
        except Exception as e:
            logging.error(f"Error copying {file_path} to {destination}: {e}")


# -------------------------
# Processing
# -------------------------
def process_file(
    file_path: Path,
    root_dir: Path,
    compiled_patterns: dict,
    output_dir: Path,
    stats: Statistics,
    metadata: MetadataStore,
) -> None:
    try:
        st = file_path.stat()
        mtime = st.st_mtime
        fsize = st.st_size
    except Exception as e:
        logging.error(f"Cannot access file stats for {file_path}: {e}")
        return

    db_key = str(file_path.resolve())
    stored = metadata.get(db_key)
    if stored and stored == (mtime, fsize):
        logging.info(f"Skipping unchanged file: {file_path}")
        return

    logging.info(f"Processing file: {file_path}")
    text_content = extract_text(file_path)
    rel_path = safe_relpath(file_path, root_dir)

    if text_content:
        found = search_keywords(text_content, compiled_patterns)
        if found:
            for group, keywords in found.items():
                stats.update_keyword_group(group)
                for keyword in keywords:
                    stats.update_keyword(group, keyword)
                    stats.update_keyword_file_path(keyword, rel_path)
                stats.record_matched_file(rel_path, group, fsize)

            if len(found) > 1:
                stats.increment_multiple_keywords()

            copy_file_to_groups(file_path, rel_path, found.keys(), output_dir)

            for group, keywords in found.items():
                for keyword in keywords:
                    logging.info(f"Found keyword '{keyword}' in {file_path} for group '{group}'")

            stats.update_filetype(file_path.suffix.lower())
    else:
        logging.warning(f"No text extracted from {file_path}")

    metadata.set(db_key, mtime, fsize)


def generate_group_statistics(stats: Statistics, output_dir: Path) -> None:
    all_groups = list(stats.keyword_group_counts.keys())
    for group in all_groups:
        group_dict = stats.group_to_dict(group)
        group_file = output_dir / f"statistics_{group}.json"
        try:
            with open(group_file, "w", encoding="utf-8") as f:
                json.dump(group_dict, f, indent=4)
            logging.info(f"Group statistics written to {group_file}")
        except Exception as e:
            logging.error(f"Error writing group statistics for {group} to {group_file}: {e}")


def generate_statistics(stats: Statistics, output_dir: Path) -> None:
    stats_file = output_dir / "statistics.json"
    try:
        with open(stats_file, "w", encoding="utf-8") as f:
            json.dump(stats.to_dict(), f, indent=4)
        logging.info(f"Statistics written to {stats_file}")
    except Exception as e:
        logging.error(f"Error writing statistics to {stats_file}: {e}")

    generate_group_statistics(stats, output_dir)


def consumer(
    queue: Queue,
    root_dir: Path,
    compiled_patterns: dict,
    output_dir: Path,
    stats: Statistics,
    metadata: MetadataStore,
) -> None:
    while True:
        item = queue.get()
        if item is None:
            queue.task_done()
            break
        try:
            process_file(item, root_dir, compiled_patterns, output_dir, stats, metadata)
        finally:
            queue.task_done()


def iter_files(root_dir: Path):
    """
    Iterates all files under root_dir. Preserves original behavior:
    - includes files with no extension (treated as .txt by extract_text)
    """
    for dirpath, _, filenames in os.walk(root_dir):
        for fn in filenames:
            yield Path(dirpath) / fn


# -------------------------
# Main
# -------------------------
def main():
    parser = argparse.ArgumentParser(description="OfficeScrap - keyword search and collect files by group.")
    parser.add_argument("--data", required=True, help="Root folder to analyze.")
    parser.add_argument("--config", required=True, help="YAML configuration with Keyword_Groups.")
    parser.add_argument("--output", required=True, help="Output folder for grouped copies and statistics.")
    parser.add_argument("--threads", type=int, default=8, help="Number of worker threads.")
    parser.add_argument("--log", default="officescrap.log", help="Log file name (written inside output folder).")
    parser.add_argument("--verbose", action="store_true", help="Verbose console logging.")
    args = parser.parse_args()

    root_dir = Path(args.data).expanduser().resolve()
    output_dir = Path(args.output).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    log_path = output_dir / args.log
    setup_logging(log_path, args.verbose)

    config_path = Path(args.config).expanduser().resolve()
    keyword_groups = load_config(config_path)
    compiled_patterns = compile_keyword_patterns(keyword_groups)

    # Metadata DB inside output folder to avoid scattering artifacts
    metadata_db_path = output_dir / "metadata_store.sqlite"
    metadata = MetadataStore(metadata_db_path)

    stats = Statistics()
    q = Queue(maxsize=10_000)

    threads = []
    for _ in range(max(1, args.threads)):
        t = Thread(
            target=consumer,
            args=(q, root_dir, compiled_patterns, output_dir, stats, metadata),
            daemon=True,
        )
        t.start()
        threads.append(t)

    # Producer
    for fp in iter_files(root_dir):
        q.put(fp)

    # Stop signals
    for _ in threads:
        q.put(None)

    q.join()
    for t in threads:
        t.join()

    generate_statistics(stats, output_dir)
    metadata.close()

    logging.info("Done.")


if __name__ == "__main__":
    main()
